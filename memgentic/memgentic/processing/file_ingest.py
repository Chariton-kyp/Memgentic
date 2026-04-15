"""File and URL content extraction for manual uploads."""

from __future__ import annotations

import re

import structlog

logger = structlog.get_logger()


def extract_text_from_file(content: bytes, mime_type: str) -> str:
    """Extract plain text from uploaded file content.

    Supported MIME types:
    - text/plain, text/markdown, text/csv, text/html — decoded as UTF-8
    - application/pdf — extracted via pdfminer.six (optional dep)
    - application/json — decoded as UTF-8
    - application/vnd.openxmlformats-officedocument.wordprocessingml.document (.docx)
      — extracted via python-docx (optional dep)
    - application/rtf, text/rtf (.rtf) — extracted via striprtf (optional dep)
    - application/vnd.openxmlformats-officedocument.spreadsheetml.sheet (.xlsx)
      — extracted via openpyxl (optional dep)
    - application/vnd.openxmlformats-officedocument.presentationml.presentation (.pptx)
      — extracted via python-pptx (optional dep)
    - application/epub+zip (.epub) — extracted via ebooklib (optional dep)

    Returns the extracted text content.

    Raises:
        ValueError: If the MIME type is unsupported or extraction fails.
    """
    text_types = {
        "text/plain",
        "text/markdown",
        "text/csv",
        "text/html",
        "text/x-markdown",
        "application/json",
        "application/xml",
        "text/xml",
    }

    if mime_type in text_types:
        text = content.decode("utf-8", errors="replace")
        # Strip HTML tags for text/html
        if mime_type == "text/html":
            text = _strip_html(text)
        return text.strip()

    if mime_type == "application/pdf":
        return _extract_pdf(content)

    # Microsoft Word (.docx)
    docx_type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    if mime_type == docx_type:
        return _extract_docx(content)

    # Rich Text Format (.rtf)
    if mime_type in ("application/rtf", "text/rtf"):
        return _extract_rtf(content)

    # Excel (.xlsx)
    xlsx_type = "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"
    if mime_type == xlsx_type:
        return _extract_xlsx(content)

    # PowerPoint (.pptx)
    pptx_type = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    if mime_type == pptx_type:
        return _extract_pptx(content)

    # EPUB (.epub)
    if mime_type == "application/epub+zip":
        return _extract_epub(content)

    raise ValueError(f"Unsupported file type: {mime_type}")


def _extract_pdf(content: bytes) -> str:
    """Extract text from PDF bytes using pdfminer.six."""
    try:
        from io import BytesIO

        from pdfminer.high_level import extract_text as pdf_extract_text

        text = pdf_extract_text(BytesIO(content))
        return text.strip()
    except ImportError:
        logger.warning(
            "file_ingest.pdfminer_not_installed",
            msg="pdfminer.six is not installed — PDF extraction unavailable. "
            "Install with: pip install pdfminer.six",
        )
        raise ValueError(
            "PDF extraction requires pdfminer.six. Install with: pip install pdfminer.six"
        ) from None
    except Exception as exc:
        logger.error("file_ingest.pdf_extraction_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from PDF: {exc}") from exc


def _strip_html(html: str) -> str:
    """Remove HTML tags and decode entities, returning plain text."""
    # Remove script and style blocks
    html = re.sub(r"<(script|style)[^>]*>.*?</\1>", "", html, flags=re.DOTALL | re.IGNORECASE)
    # Remove HTML tags
    text = re.sub(r"<[^>]+>", " ", html)
    # Decode common HTML entities
    text = text.replace("&amp;", "&")
    text = text.replace("&lt;", "<")
    text = text.replace("&gt;", ">")
    text = text.replace("&quot;", '"')
    text = text.replace("&#39;", "'")
    text = text.replace("&nbsp;", " ")
    # Collapse whitespace
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _extract_docx(content: bytes) -> str:
    """Extract text from .docx (Word) bytes using python-docx."""
    try:
        from io import BytesIO

        from docx import Document

        doc = Document(BytesIO(content))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        text = "\n\n".join(paragraphs)
        return text.strip()
    except ImportError:
        logger.warning(
            "file_ingest.python_docx_not_installed",
            msg="python-docx is not installed — .docx extraction unavailable. "
            "Install with: pip install python-docx",
        )
        raise ValueError(
            "DOCX extraction requires python-docx. Install with: pip install python-docx"
        ) from None
    except Exception as exc:
        logger.error("file_ingest.docx_extraction_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from DOCX: {exc}") from exc


def _extract_rtf(content: bytes) -> str:
    """Extract text from .rtf bytes using striprtf."""
    try:
        from striprtf.striprtf import rtf_to_text

        raw = content.decode("utf-8", errors="replace")
        text = rtf_to_text(raw)
        return text.strip()
    except ImportError:
        logger.warning(
            "file_ingest.striprtf_not_installed",
            msg="striprtf is not installed — .rtf extraction unavailable. "
            "Install with: pip install striprtf",
        )
        raise ValueError(
            "RTF extraction requires striprtf. Install with: pip install striprtf"
        ) from None
    except Exception as exc:
        logger.error("file_ingest.rtf_extraction_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from RTF: {exc}") from exc


def _extract_xlsx(content: bytes) -> str:
    """Extract text from .xlsx (Excel) bytes using openpyxl."""
    try:
        from io import BytesIO

        from openpyxl import load_workbook

        wb = load_workbook(BytesIO(content), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet in wb.sheetnames:
            ws = wb[sheet]
            rows.append(f"## Sheet: {sheet}\n")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    rows.append(" | ".join(cells))
        wb.close()
        return "\n".join(rows).strip()
    except ImportError:
        logger.warning(
            "file_ingest.openpyxl_not_installed",
            msg="openpyxl is not installed — .xlsx extraction unavailable. "
            "Install with: pip install openpyxl",
        )
        raise ValueError(
            "XLSX extraction requires openpyxl. Install with: pip install openpyxl"
        ) from None
    except Exception as exc:
        logger.error("file_ingest.xlsx_extraction_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from XLSX: {exc}") from exc


def _extract_pptx(content: bytes) -> str:
    """Extract text from .pptx (PowerPoint) bytes using python-pptx."""
    try:
        from io import BytesIO

        from pptx import Presentation

        prs = Presentation(BytesIO(content))
        slides_text: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts: list[str] = [f"## Slide {i}"]
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            slides_text.append("\n".join(parts))
        return "\n\n".join(slides_text).strip()
    except ImportError:
        logger.warning(
            "file_ingest.python_pptx_not_installed",
            msg="python-pptx is not installed — .pptx extraction unavailable. "
            "Install with: pip install python-pptx",
        )
        raise ValueError(
            "PPTX extraction requires python-pptx. Install with: pip install python-pptx"
        ) from None
    except Exception as exc:
        logger.error("file_ingest.pptx_extraction_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from PPTX: {exc}") from exc


def _extract_epub(content: bytes) -> str:
    """Extract text from .epub bytes using ebooklib + html stripping."""
    try:
        from io import BytesIO

        import ebooklib
        from ebooklib import epub

        book = epub.read_epub(BytesIO(content))
        chapters: list[str] = []
        for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
            html = item.get_content().decode("utf-8", errors="replace")
            text = _strip_html(html)
            if text.strip():
                chapters.append(text.strip())
        return "\n\n".join(chapters).strip()
    except ImportError:
        logger.warning(
            "file_ingest.ebooklib_not_installed",
            msg="ebooklib is not installed — .epub extraction unavailable. "
            "Install with: pip install ebooklib",
        )
        raise ValueError(
            "EPUB extraction requires ebooklib. Install with: pip install ebooklib"
        ) from None
    except Exception as exc:
        logger.error("file_ingest.epub_extraction_failed", error=str(exc))
        raise ValueError(f"Failed to extract text from EPUB: {exc}") from exc


async def extract_text_from_url(url: str) -> tuple[str, str]:
    """Fetch a URL and extract text content.

    Returns:
        A tuple of (extracted_text, page_title).

    Raises:
        ValueError: If the URL cannot be fetched or content extracted.
    """
    try:
        import httpx
    except ImportError:
        raise ValueError("URL import requires httpx. Install with: pip install httpx") from None

    try:
        async with httpx.AsyncClient(
            follow_redirects=True,
            timeout=30.0,
            headers={
                "User-Agent": "Memgentic/1.0 (memory indexer)",
                "Accept": "text/html,application/xhtml+xml,text/plain,*/*",
            },
        ) as client:
            response = await client.get(url)
            response.raise_for_status()
    except httpx.HTTPStatusError as exc:
        raise ValueError(f"HTTP error fetching URL: {exc.response.status_code}") from exc
    except httpx.RequestError as exc:
        raise ValueError(f"Failed to fetch URL: {exc}") from exc

    content_type = response.headers.get("content-type", "text/html")
    raw = response.text

    # Extract title from HTML
    title = _extract_html_title(raw) if "html" in content_type else ""

    # Strip HTML to get plain text
    text = _strip_html(raw) if "html" in content_type or "xml" in content_type else raw

    if not text.strip():
        raise ValueError("No text content could be extracted from the URL")

    # Truncate very long pages to a reasonable size
    max_chars = 50000
    if len(text) > max_chars:
        text = text[:max_chars] + "\n\n[Content truncated]"

    return text.strip(), title


def _extract_html_title(html: str) -> str:
    """Extract the <title> from HTML content."""
    match = re.search(r"<title[^>]*>(.*?)</title>", html, re.IGNORECASE | re.DOTALL)
    if match:
        title = match.group(1).strip()
        # Clean up whitespace and entities
        title = re.sub(r"\s+", " ", title)
        title = title.replace("&amp;", "&")
        title = title.replace("&lt;", "<")
        title = title.replace("&gt;", ">")
        title = title.replace("&quot;", '"')
        return title
    return ""
